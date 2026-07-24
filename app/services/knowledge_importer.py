from __future__ import annotations

import csv
import hashlib
import io
import json
import re
import zipfile
from dataclasses import asdict, dataclass
from datetime import UTC, datetime
from html.parser import HTMLParser
from pathlib import Path, PurePosixPath
from typing import Callable

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation


MAX_MEMBER_BYTES = 25 * 1024 * 1024
MAX_TOTAL_BYTES = 250 * 1024 * 1024


@dataclass(slots=True)
class InventoryItem:
    path: str
    archive_path: str
    extension: str
    size_bytes: int
    category: str
    importable: bool
    reason: str


class _TextHTMLParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []
        self.skip_depth = 0

    def handle_starttag(
        self,
        tag: str,
        attrs: list[tuple[str, str | None]],
    ) -> None:
        if tag in {"script", "style", "noscript", "svg"}:
            self.skip_depth += 1
            return
        if not self.skip_depth and tag in {
            "p",
            "div",
            "li",
            "h1",
            "h2",
            "h3",
            "br",
            "tr",
            "td",
        }:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in {"script", "style", "noscript", "svg"}:
            self.skip_depth = max(0, self.skip_depth - 1)
            return
        if not self.skip_depth and tag in {
            "p",
            "div",
            "li",
            "h1",
            "h2",
            "h3",
            "tr",
        }:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if self.skip_depth:
            return
        value = data.strip()
        if value:
            self.parts.append(value)


def _decode(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise ValueError("unsupported text encoding")


def _plain_text(data: bytes) -> str:
    return _decode(data)


def _json_text(data: bytes) -> str:
    payload = json.loads(_decode(data))
    return json.dumps(payload, ensure_ascii=False, indent=2)


def _csv_text(data: bytes) -> str:
    reader = csv.reader(io.StringIO(_decode(data)))
    return "\n".join(
        "| " + " | ".join(str(cell).replace("|", "\\|") for cell in row) + " |"
        for row in reader
    )


def _html_text(data: bytes) -> str:
    parser = _TextHTMLParser()
    parser.feed(_decode(data))
    text = "".join(parser.parts)
    return "\n".join(
        line
        for raw in text.splitlines()
        if (line := re.sub(r"\s+", " ", raw).strip())
    )


def _coze_knowledge_document(text: str) -> str:
    """Extract the actual document preview from a saved Coze knowledge page."""
    marker = "预览原始文档"
    if marker not in text:
        return ""
    content = text.split(marker, 1)[1].strip()
    return "\n".join(
        line
        for line in content.splitlines()
        if not re.fullmatch(r"\d+\s*/\s*\d+|\d+%", line.strip())
    )


def _pdf_text(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    return "\n\n".join(
        text
        for page in reader.pages
        if (text := (page.extract_text() or "").strip())
    )


def _docx_text(data: bytes) -> str:
    document = Document(io.BytesIO(data))
    paragraphs = [
        paragraph.text.strip()
        for paragraph in document.paragraphs
        if paragraph.text.strip()
    ]
    for table in document.tables:
        for row in table.rows:
            paragraphs.append(
                " | ".join(cell.text.strip() for cell in row.cells)
            )
    return "\n\n".join(paragraphs)


def _pptx_text(data: bytes) -> str:
    presentation = Presentation(io.BytesIO(data))
    parts = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
        if slide_parts:
            parts.append(f"## 第 {index} 页\n\n" + "\n\n".join(slide_parts))
    return "\n\n".join(parts)


def _xlsx_text(data: bytes) -> str:
    workbook = load_workbook(
        io.BytesIO(data),
        read_only=True,
        data_only=True,
    )
    parts = []
    for sheet in workbook.worksheets:
        parts.append(f"## {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(values):
                parts.append(
                    "| "
                    + " | ".join(value.replace("|", "\\|") for value in values)
                    + " |"
                )
    return "\n".join(parts)


CONVERTERS: dict[str, Callable[[bytes], str]] = {
    ".md": _plain_text,
    ".txt": _plain_text,
    ".json": _json_text,
    ".csv": _csv_text,
    ".html": _html_text,
    ".htm": _html_text,
    ".pdf": _pdf_text,
    ".docx": _docx_text,
    ".pptx": _pptx_text,
    ".xlsx": _xlsx_text,
}


class KnowledgeImporter:
    def __init__(self, zip_path: Path) -> None:
        self.zip_path = zip_path

    def inventory(self) -> list[InventoryItem]:
        total_size = 0
        items = []
        with zipfile.ZipFile(self.zip_path) as archive:
            for member in archive.infolist():
                if member.is_dir():
                    continue
                self._validate_member_path(member.filename)
                total_size += member.file_size
                if total_size > MAX_TOTAL_BYTES:
                    raise ValueError("zip uncompressed size exceeds safety limit")
                display_path = self._repair_zip_name(member.filename)
                extension = PurePosixPath(display_path).suffix.lower()
                category, reason = self._classify(display_path, extension)
                importable = (
                    extension in CONVERTERS
                    and member.file_size <= MAX_MEMBER_BYTES
                    and category != "platform_metadata"
                )
                if member.file_size > MAX_MEMBER_BYTES:
                    reason = "file exceeds 25 MiB safety limit"
                items.append(
                    InventoryItem(
                        path=display_path,
                        archive_path=member.filename,
                        extension=extension,
                        size_bytes=member.file_size,
                        category=category,
                        importable=importable,
                        reason=reason,
                    )
                )
        return items

    def import_documents(
        self,
        output_dir: Path,
        *,
        include_documents: bool = False,
        include_unclassified: bool = False,
    ) -> list[dict]:
        inventory = self.inventory()
        output_dir.mkdir(parents=True, exist_ok=True)
        imported = []
        by_path = {item.archive_path: item for item in inventory}
        with zipfile.ZipFile(self.zip_path) as archive:
            for archive_path, item in by_path.items():
                if not item.importable:
                    continue
                if item.category == "document" and not include_documents:
                    continue
                if (
                    item.category == "unclassified"
                    and not include_unclassified
                ):
                    continue
                raw = archive.read(archive_path)
                converter = CONVERTERS[item.extension]
                try:
                    content = converter(raw).strip()
                    if (
                        item.category == "knowledge"
                        and item.extension in {".html", ".htm"}
                    ):
                        content = _coze_knowledge_document(content)
                except Exception as exc:
                    imported.append(
                        {
                            "source_path": item.path,
                            "status": "failed",
                            "reason": str(exc),
                        }
                    )
                    continue
                if not content:
                    imported.append(
                        {
                            "source_path": item.path,
                            "status": "skipped",
                            "reason": "no embedded knowledge document found",
                        }
                    )
                    continue

                digest = hashlib.sha256(
                    item.path.encode("utf-8")
                ).hexdigest()[:10]
                stem = self._safe_stem(PurePosixPath(item.path).stem)
                output_path = output_dir / f"{stem}-{digest}.md"
                front_matter = (
                    "---\n"
                    f"source_archive: {json.dumps(self.zip_path.name, ensure_ascii=False)}\n"
                    f"source_path: {json.dumps(item.path, ensure_ascii=False)}\n"
                    f"imported_at: {datetime.now(UTC).isoformat()}\n"
                    "verified: false\n"
                    "---\n\n"
                )
                output_path.write_text(
                    front_matter + content + "\n",
                    encoding="utf-8",
                )
                imported.append(
                    {
                        "source_path": item.path,
                        "output_path": str(output_path),
                        "status": "imported",
                        "characters": len(content),
                    }
                )
        return imported

    @staticmethod
    def write_inventory(
        inventory: list[InventoryItem],
        output_path: Path,
    ) -> None:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(
            json.dumps(
                [asdict(item) for item in inventory],
                ensure_ascii=False,
                indent=2,
            )
            + "\n",
            encoding="utf-8",
        )

    @staticmethod
    def _validate_member_path(path: str) -> None:
        pure = PurePosixPath(path)
        if pure.is_absolute() or ".." in pure.parts:
            raise ValueError(f"unsafe zip member path: {path}")
        if "\x00" in path:
            raise ValueError("zip member contains NUL byte")

    @staticmethod
    def _classify(path: str, extension: str) -> tuple[str, str]:
        lower = path.lower()
        if (
            path.startswith("__MACOSX/")
            or "/._" in path
            or lower.endswith(".ds_store")
            or "_files/" in lower
        ):
            return "platform_metadata", "saved-page asset or macOS metadata"
        if any(
            token in lower
            for token in (
                "workflow",
                "_flow",
                "plugin",
                "bot_config",
                "prompt",
                "智能体 - 扣子",
                "智能体平台",
            )
        ):
            return "platform_metadata", "platform implementation is excluded"
        if any(
            token in lower
            for token in (
                "knowledge",
                "dataset",
                "document",
                "知识",
                "资料",
                "文档",
            )
        ):
            return "knowledge", "knowledge-like path"
        if extension in {".pdf", ".docx", ".pptx", ".xlsx", ".md", ".txt"}:
            return "document", "human-readable document"
        return "unclassified", "requires manual review"

    @staticmethod
    def _repair_zip_name(value: str) -> str:
        try:
            return value.encode("cp437").decode("utf-8")
        except (UnicodeEncodeError, UnicodeDecodeError):
            return value

    @staticmethod
    def _safe_stem(value: str) -> str:
        value = re.sub(r"[^\w\u4e00-\u9fff-]+", "-", value).strip("-")
        return value[:60] or "document"
